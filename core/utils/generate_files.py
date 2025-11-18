import asyncio
import os
import re
import io

from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.pdfgen import canvas
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from typing import Any, List, Optional, Sequence
from xml.sax.saxutils import escape

from core.models.worklet import Worklet
from core.utils.sanitize_filename import sanitize_filename
import uuid

CUSTOM_PAGE_SIZE = (
    750,
    900,
)  # Width x Height in points (1 point = 1/72 inch) used by pdf

# small vertical gap (in inches) used between dynamic blocks in PPT
gap = 0.3
DEFAULT_PPT_GAP_INCH = 0.3


# ---------------------------
# Utility helpers
# ---------------------------


def ensure_list(value: Any) -> List:
    """Return a list for list-like inputs, or [] for None, or [value] for single non-list values."""
    if value is None:
        return []
    if isinstance(value, list):
        return value
    if isinstance(value, (tuple, set)):
        return list(value)
    # if it's a string, treat as single element list (but prefer not to split)
    return [value]


_BULLET_PREFIX_RE = re.compile(r"^\s*[\u2022\-*\u2023]+\s*")


def _strip_bullet_prefix(text: str) -> str:
    return _BULLET_PREFIX_RE.sub("", text)


def normalize_text_list(value: Any, *, split_on_delimiters: bool = True) -> List[str]:
    """Convert raw values into a list of trimmed, non-empty strings."""

    candidates = ensure_list(value)
    normalized: List[str] = []
    for candidate in candidates:
        if candidate is None:
            continue
        if isinstance(candidate, str):
            text = candidate.replace("\r\n", "\n").strip()
            if not text:
                continue
            if split_on_delimiters:
                segments = re.split(r"[\r\n]+|\s*[;\u2022]\s*", text)
                for segment in segments:
                    cleaned = _strip_bullet_prefix(segment.strip())
                    if cleaned:
                        normalized.append(cleaned)
            else:
                cleaned = _strip_bullet_prefix(text)
                if cleaned:
                    normalized.append(cleaned)
        else:
            text = str(candidate).strip()
            if text:
                normalized.append(text)
    return normalized


def format_multiline_pdf_bullet(text: str) -> str:
    """Render a single bullet with optional indented sub-lines for PDF output."""

    if not text:
        return ""
    lines = [
        escape(line.strip())
        for line in text.replace("\r\n", "\n").split("\n")
        if line and line.strip()
    ]
    if not lines:
        return ""
    if len(lines) == 1:
        return f"• {lines[0]}"
    first, *rest = lines
    rest_markup = "<br/>".join(f"&nbsp;&nbsp;{line}" for line in rest)
    return f"• {first}<br/>{rest_markup}"


def format_multiline_ppt_bullet(text: str) -> str:
    """Render a single bullet with optional indented sub-lines for PPT output."""

    if not text:
        return ""
    lines = [
        line.strip()
        for line in text.replace("\r\n", "\n").split("\n")
        if line and line.strip()
    ]
    if not lines:
        return ""
    if len(lines) == 1:
        return f"• {lines[0]}"
    first, *rest = lines
    continuation = "\n   ".join(rest)
    return f"• {first}\n   {continuation}"


def safe_get(obj: Any, keys: Sequence[str], default: Any = None) -> Any:
    """
    Safely retrieve the first non-empty value from obj for any of the provided keys.
    Supports both dicts and objects with attributes.
    """
    if obj is None:
        return default
    for key in keys:
        if isinstance(obj, dict):
            if key in obj and obj[key] not in (None, ""):
                return obj[key]
        else:
            if hasattr(obj, key):
                val = getattr(obj, key)
                if val not in (None, ""):
                    return val
    return default


def extract_reference_field(ref: Any, possible_keys: Sequence[str]) -> Optional[str]:
    """Extract a string field from a Reference object or a dict with multiple possible key names."""
    return safe_get(ref, possible_keys, default=None)


# Mapping of canonical fields to possible keys/names in older dicts or new models
FIELD_KEYS = {
    "title": ["title", "Title"],
    "problem_statement": [
        "problem_statement",
        "Problem Statement",
        "problem statement",
        "problemStatement",
    ],
    "description": ["description", "Description"],
    "challenge_use_case": [
        "challenge_use_case",
        "Challenge / Use Case",
        "Challenge",
        "challenge Use Case",
    ],
    "deliverables": ["deliverables", "Deliverables"],
    "kpis": ["kpis", "KPIs", "KPI", "Kpis"],
    "prerequisites": ["prerequisites", "Prerequisites"],
    "infrastructure_requirements": [
        "infrastructure_requirements",
        "Infrastructure Requirements",
        "Infrastructure",
    ],
    "tech_stack": [
        "tech_stack",
        "Tentative Tech Stack",
        "Tentative_Tech_Stack",
        "Tentative Tech Stack",
    ],
    "milestones": [
        "milestones",
        "Milestones (6 months)",
        "Milestones",
        "milestones_6_months",
    ],
    "references": ["references", "Reference Work", "Reference", "Reference_Work"],
}


# ---------------------------
# Main generator
# ---------------------------
async def generate_file(worklet: Worklet, thread_id: str) -> None:
    """
    Entrypoint to generate both PDF and PPTX for a worklet.
    Accepts either a Pydantic Worklet instance or a legacy dict-like worklet.
    """
    ppt_path = os.path.join("data/threads", thread_id, "generated_worklets/pptx")
    pdf_path = os.path.join("data/threads", thread_id, "generated_worklets/pdf")
    os.makedirs(ppt_path, exist_ok=True)
    os.makedirs(pdf_path, exist_ok=True)

    # Resolve title safely (handle dicts or objects)
    title = safe_get(worklet, FIELD_KEYS["title"], default="untitled")
    safe_title = sanitize_filename(title)

    filename = safe_title if safe_title else f"untitled_{uuid.uuid4().hex[:8]}"
    filename_pdf = os.path.join(pdf_path, f"{safe_title}.pdf")
    filename_ppt = os.path.join(ppt_path, f"{safe_title}.pptx")

    # Ensure directories exist
    os.makedirs(os.path.dirname(filename_pdf) or ".", exist_ok=True)
    os.makedirs(os.path.dirname(filename_ppt) or ".", exist_ok=True)

    loop = asyncio.get_running_loop()
    await loop.run_in_executor(None, create_pdf, filename_pdf, worklet)
    await loop.run_in_executor(None, create_ppt, filename_ppt, worklet)


# ---------------------------
# PDF CREATION
# ---------------------------
def create_pdf(filename: str, worklet: Worklet, in_memory: bool = False):
    """
    Create a PDF summarizing the worklet.
    This function is resilient to missing fields and supports both dict-like and attribute-like worklets.
    """
    try:
        # Use SimpleDocTemplate and platypus flowables so ReportLab handles pagination
        width, height = CUSTOM_PAGE_SIZE

        styles = getSampleStyleSheet()
        header_style = ParagraphStyle(
            "header_style",
            parent=styles["Heading1"],
            fontSize=20,
            textColor=colors.darkblue,
        )
        normal_style = ParagraphStyle(
            "normal_style", parent=styles["BodyText"], fontSize=12, leading=15
        )
        bullet_style = ParagraphStyle(
            "bullet_style",
            parent=styles["BodyText"],
            fontSize=12,
            leftIndent=20,
            bulletIndent=10,
        )

        elements = []

        # Title & core fields - only add if present and non-empty
        title = safe_get(worklet, FIELD_KEYS["title"])
        if title:
            elements.append(Paragraph(f"<b>Title:</b> {title}", header_style))
            elements.append(Spacer(1, 6))

        problem = safe_get(worklet, FIELD_KEYS["problem_statement"])
        if problem:
            elements.append(
                Paragraph(f"<b>Problem Statement:</b> {problem}", normal_style)
            )
            elements.append(Spacer(1, 6))

        desc = safe_get(worklet, FIELD_KEYS["description"])
        if desc:
            elements.append(Paragraph(f"<b>Description:</b> {desc}", normal_style))
            elements.append(Spacer(1, 6))

        challenge = safe_get(worklet, FIELD_KEYS["challenge_use_case"])
        if challenge:
            elements.append(
                Paragraph(f"<b>Challenge / Use Case:</b> {challenge}", normal_style)
            )
            elements.append(Spacer(1, 6))

        raw_deliverables = safe_get(worklet, FIELD_KEYS["deliverables"])
        deliverables = normalize_text_list(raw_deliverables)
        if deliverables:
            elements.append(Paragraph("<b>Deliverables:</b>", normal_style))
            for item in deliverables:
                elements.append(Paragraph(f"• {item}", bullet_style))
            elements.append(Spacer(1, 6))

        # KPIs (list)
        raw_kpis = safe_get(worklet, FIELD_KEYS["kpis"])
        kpis = normalize_text_list(raw_kpis, split_on_delimiters=False)
        if kpis:
            elements.append(Paragraph("<b>KPIs:</b>", normal_style))
            for kpi in kpis:
                bullet_text = format_multiline_pdf_bullet(kpi)
                if bullet_text:
                    elements.append(Paragraph(bullet_text, bullet_style))
            elements.append(Spacer(1, 6))

        # Prerequisites (list)
        raw_prereqs = safe_get(worklet, FIELD_KEYS["prerequisites"])
        prereqs = normalize_text_list(raw_prereqs)
        if prereqs:
            elements.append(Paragraph("<b>Prerequisites:</b>", normal_style))
            for prereq in prereqs:
                elements.append(Paragraph(f"• {prereq}", bullet_style))
            elements.append(Spacer(1, 6))

        # Infra & Tech Stack
        infra = safe_get(worklet, FIELD_KEYS["infrastructure_requirements"])
        if infra:
            elements.append(
                Paragraph(f"<b>Infrastructure Requirements:</b> {infra}", normal_style)
            )
            elements.append(Spacer(1, 6))

        tech = safe_get(worklet, FIELD_KEYS["tech_stack"])
        if tech:
            elements.append(
                Paragraph(f"<b>Tentative Tech Stack:</b> {tech}", normal_style)
            )
            elements.append(Spacer(1, 6))

        # Milestones (dict)
        milestones = safe_get(worklet, FIELD_KEYS["milestones"])
        if isinstance(milestones, dict) and milestones:
            elements.append(Paragraph("<b>Milestones (6 months):</b>", normal_style))
            # Prefer M2/M4/M6 ordering if present
            for key in ("M2", "M4", "M6"):
                if key in milestones and milestones[key] not in (None, ""):
                    elements.append(
                        Paragraph(f"• {key}: {milestones[key]}", bullet_style)
                    )
            # Add any other milestones
            for k, v in milestones.items():
                if k not in ("M2", "M4", "M6") and v not in (None, ""):
                    elements.append(Paragraph(f"• {k}: {v}", bullet_style))
            elements.append(Spacer(1, 6))

        # References: support list of Reference objects or dicts
        raw_refs = safe_get(worklet, FIELD_KEYS["references"]) or []
        refs = ensure_list(raw_refs)
        if refs:
            elements.append(Paragraph("<b>References:</b>", normal_style))
            for ref in refs:
                # Support both dict and object forms
                title_r = extract_reference_field(ref, ["title", "Title"])
                link_r = extract_reference_field(ref, ["link", "Link", "url", "URL"])
                desc_r = extract_reference_field(
                    ref, ["description", "Description", "abstract"]
                )
                tag_r = extract_reference_field(ref, ["tag", "Tag", "source"])

                # Create reference bullet point with plain title and clickable "link"
                if title_r and link_r:
                    # Plain title followed by clickable "link" word
                    composed = (
                        f'• {title_r} <a href="{link_r}" color="blue"><u>link</u></a>'
                    )
                elif title_r:
                    # Title without link
                    composed = f"• {title_r}"
                elif link_r:
                    # Link without title - show clickable "link" word
                    composed = f'• <a href="{link_r}" color="blue"><u>link</u></a>'
                else:
                    # Fallback to stringifying the reference
                    composed = f"• {str(ref)}"

                if composed:
                    elements.append(Paragraph(composed, bullet_style))
            elements.append(Spacer(1, 6))

        # If no elements were added, add a minimal notice so the PDF is not blank
        if not elements:
            elements.append(
                Paragraph("No content available for this worklet.", normal_style)
            )

        # Build document - supports pagination automatically
        if in_memory:
            buffer = io.BytesIO()
            doc = SimpleDocTemplate(
                buffer,
                pagesize=CUSTOM_PAGE_SIZE,
                leftMargin=40,
                rightMargin=40,
                topMargin=60,
                bottomMargin=40,
            )
            doc.build(elements)
            data = buffer.getvalue()
            buffer.close()
            return data
        else:
            doc = SimpleDocTemplate(
                filename,
                pagesize=CUSTOM_PAGE_SIZE,
                leftMargin=40,
                rightMargin=40,
                topMargin=60,
                bottomMargin=40,
            )
            doc.build(elements)

    except Exception as e:
        print(f"Failed to generate PDF {filename}: {e}")


# ---------------------------
# PPT CREATION
# ---------------------------
def create_ppt(output_filename: str, worklet: Worklet, in_memory: bool = False):
    """
    Create a single-slide PPTX summarizing the worklet.
    Robust to missing fields and supports dict-like and attribute-like worklets.
    """
    try:
        prs = Presentation()
        prs.slide_width = Pt(750)
        prs.slide_height = Pt(1100)
        slide_layout = prs.slide_layouts[6]

        # Pagination state (top measured in inches)
        slide = prs.slides.add_slide(slide_layout)
        top = 0.5  # inches from top (top margin)
        gap = DEFAULT_PPT_GAP_INCH
        top_margin = 0.5
        bottom_margin = 0.5
        slide_height_in = prs.slide_height.inches

        # helper to ensure there is space on current slide, otherwise create a new slide
        def _ensure_space(needed_height: float):
            nonlocal slide, top
            usable_bottom = slide_height_in - bottom_margin
            # If content would overflow the usable area, start a new slide
            if top + needed_height > usable_bottom:
                slide = prs.slides.add_slide(slide_layout)
                top = top_margin

        # Helper to safely retrieve textual fields
        def _text_for(keys: Sequence[str]) -> str:
            val = safe_get(worklet, keys)
            return str(val).strip() if val not in (None, "") else ""

        # Title block
        title_text = _text_for(FIELD_KEYS["title"])
        if title_text:
            # compute estimated height and ensure page space
            try:
                est_h = estimate_height_wrapped_Title(title_text)
                _ensure_space(est_h + gap)
                top = add_textbox_Title(slide, "Title", title_text, top) + gap
            except NameError:
                # Fallback: create a basic textbox if helper not present
                est_h = 0.6
                _ensure_space(est_h + gap)
                left = Inches(0.5)
                top_inch = Inches(top)
                width = Inches(9.5)
                height = Inches(est_h)
                tb = slide.shapes.add_textbox(left, top_inch, width, height)
                tf = tb.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.text = title_text
                p.font.size = Pt(22)
                top += est_h + gap

        # Core textual fields (add if present)
        for field_key in (
            "problem_statement",
            "description",
            "challenge_use_case",
        ):
            text = _text_for(FIELD_KEYS[field_key])
            if text:
                # estimate height and ensure page space
                try:
                    est_h = estimate_height_wrapped_content(text)
                    _ensure_space(est_h + gap)
                    top = add_textbox(
                        slide, field_key.replace("_", " ").title(), text, top
                    )
                except NameError:
                    # Fallback simple textbox
                    est_h = 0.8
                    _ensure_space(est_h + gap)
                    left = Inches(0.5)
                    top_inch = Inches(top)
                    width = Inches(9.5)
                    height = Inches(est_h)
                    tb = slide.shapes.add_textbox(left, top_inch, width, height)
                    tf = tb.text_frame
                    tf.clear()
                    p = tf.paragraphs[0]
                    p.text = f"{field_key.replace('_', ' ').title()}: {text}"
                    p.font.size = Pt(14)
                    top += est_h + gap

        # Deliverables (list)
        raw_deliverables = safe_get(worklet, FIELD_KEYS["deliverables"])
        deliverables = normalize_text_list(raw_deliverables)
        if deliverables:
            deliverables_text = "\n".join([f"• {item}" for item in deliverables])
            if deliverables_text:
                try:
                    est_h = estimate_height_wrapped_content(deliverables_text)
                    _ensure_space(est_h + gap)
                    top = add_textbox(slide, "Deliverables", deliverables_text, top)
                except NameError:
                    est_h = min(2.0, 0.3 * len(deliverables) + 0.2)
                    _ensure_space(est_h + gap)
                    left = Inches(0.5)
                    top_inch = Inches(top)
                    width = Inches(9.5)
                    height = Inches(est_h)
                    tb = slide.shapes.add_textbox(left, top_inch, width, height)
                    tf = tb.text_frame
                    tf.clear()
                    for index, item in enumerate(deliverables):
                        paragraph = (
                            tf.paragraphs[0] if index == 0 else tf.add_paragraph()
                        )
                        paragraph.text = f"• {item}"
                    top += est_h + gap

        # KPIs (list)
        raw_kpis = safe_get(worklet, FIELD_KEYS["kpis"])
        kpis = normalize_text_list(raw_kpis, split_on_delimiters=False)
        if kpis:
            formatted_kpis = [format_multiline_ppt_bullet(k) for k in kpis]
            filtered_kpis = [entry for entry in formatted_kpis if entry]
            kpi_text = "\n".join(filtered_kpis)
            if kpi_text:
                try:
                    est_h = estimate_height_wrapped_content(kpi_text)
                    _ensure_space(est_h + gap)
                    top = add_textbox(slide, "KPIs", kpi_text, top)
                except NameError:
                    est_h = min(2.0, 0.3 * len(kpi_text.split("\n")) + 0.2)
                    _ensure_space(est_h + gap)
                    left = Inches(0.5)
                    top_inch = Inches(top)
                    width = Inches(9.5)
                    height = Inches(est_h)
                    tb = slide.shapes.add_textbox(left, top_inch, width, height)
                    tf = tb.text_frame
                    tf.clear()
                    added_any = False
                    for entry in filtered_kpis:
                        if not added_any:
                            paragraph = tf.paragraphs[0]
                            added_any = True
                        else:
                            paragraph = tf.add_paragraph()
                        paragraph.text = entry
                    top += est_h + gap

        # Prerequisites (list)
        raw_prereqs = safe_get(worklet, FIELD_KEYS["prerequisites"])
        prereqs = normalize_text_list(raw_prereqs)
        if prereqs:
            preq_text = "\n".join([f"• {p}" for p in prereqs])
            if preq_text:
                try:
                    est_h = estimate_height_wrapped_content(preq_text)
                    _ensure_space(est_h + gap)
                    top = add_textbox(slide, "Prerequisites", preq_text, top)
                except NameError:
                    est_h = min(2.0, 0.3 * len(prereqs) + 0.2)
                    _ensure_space(est_h + gap)
                    left = Inches(0.5)
                    top_inch = Inches(top)
                    width = Inches(9.5)
                    height = Inches(est_h)
                    tb = slide.shapes.add_textbox(left, top_inch, width, height)
                    tf = tb.text_frame
                    tf.clear()
                    for i, p in enumerate(prereqs):
                        if i == 0:
                            p0 = tf.paragraphs[0]
                            p0.text = f"• {p}"
                        else:
                            p_par = tf.add_paragraph()
                            p_par.text = f"• {p}"
                    top += est_h + gap

        # Infra & Tech Stack
        infra = _text_for(FIELD_KEYS["infrastructure_requirements"])
        if infra:
            try:
                est_h = estimate_height_wrapped_content(infra)
                _ensure_space(est_h + gap)
                top = add_textbox(slide, "Infrastructure Requirements", infra, top)
            except NameError:
                est_h = 0.6
                _ensure_space(est_h + gap)
                left = Inches(0.5)
                top_inch = Inches(top)
                width = Inches(9.5)
                height = Inches(est_h)
                tb = slide.shapes.add_textbox(left, top_inch, width, height)
                tf = tb.text_frame
                tf.clear()
                tf.paragraphs[0].text = f"Infrastructure Requirements: {infra}"
                top += est_h + gap

        tech = _text_for(FIELD_KEYS["tech_stack"])
        if tech:
            try:
                est_h = estimate_height_wrapped_content(tech)
                _ensure_space(est_h + gap)
                top = add_textbox(slide, "Tentative Tech Stack", tech, top)
            except NameError:
                est_h = 0.6
                _ensure_space(est_h + gap)
                left = Inches(0.5)
                top_inch = Inches(top)
                width = Inches(9.5)
                height = Inches(est_h)
                tb = slide.shapes.add_textbox(left, top_inch, width, height)
                tf = tb.text_frame
                tf.clear()
                tf.paragraphs[0].text = f"Tentative Tech Stack: {tech}"
                top += est_h + gap

        # Milestones
        milestones = safe_get(worklet, FIELD_KEYS["milestones"])
        if isinstance(milestones, dict) and milestones:
            milestone_text = "\n".join(
                [f"{k}: {v}" for k, v in milestones.items() if v not in (None, "")]
            )
            if milestone_text:
                try:
                    est_h = estimate_height_wrapped_content(milestone_text)
                    _ensure_space(est_h + gap)
                    top = add_textbox(
                        slide, "Milestones (6 months)", milestone_text, top
                    )
                except NameError:
                    est_h = min(2.5, 0.25 * len(milestones) + 0.2)
                    _ensure_space(est_h + gap)
                    left = Inches(0.5)
                    top_inch = Inches(top)
                    width = Inches(9.5)
                    height = Inches(est_h)
                    tb = slide.shapes.add_textbox(left, top_inch, width, height)
                    tf = tb.text_frame
                    tf.clear()
                    for i, (k, v) in enumerate(milestones.items()):
                        if i == 0:
                            tf.paragraphs[0].text = f"{k}: {v}"
                        else:
                            p = tf.add_paragraph()
                            p.text = f"{k}: {v}"
                    top += est_h + gap

        # References block (list)
        raw_refs = safe_get(worklet, FIELD_KEYS["references"]) or []
        refs = ensure_list(raw_refs)
        if refs:
            # estimate references block height (each ref ~0.4 inch)
            est_h = min(3.5, 0.4 * len(refs) + 0.2)
            _ensure_space(est_h + gap)
            left = Inches(0.5)
            top_inch = Inches(top)
            width = Inches(9.5)
            # each ref ~0.4 inch height estimate
            height = Inches(est_h)
            textbox = slide.shapes.add_textbox(left, top_inch, width, height)
            tf = textbox.text_frame
            tf.word_wrap = True
            tf.clear()

            # Title for references
            title_para = tf.paragraphs[0]
            title_run = title_para.add_run()
            title_run.font.size = Pt(16)
            title_run.font.bold = True
            title_run.font.name = "Calibri"
            title_run.font.color.rgb = RGBColor(0x00, 0x66, 0xCC)
            title_run.text = "References:"

            for ref in refs:
                r_title = extract_reference_field(ref, ["title", "Title"])
                r_link = extract_reference_field(ref, ["link", "Link", "url", "URL"])
                r_desc = extract_reference_field(
                    ref, ["description", "Description", "abstract"]
                )
                r_tag = extract_reference_field(ref, ["tag", "Tag", "source"])

                p = tf.add_paragraph()
                p.level = 1

                # Add bullet point
                bullet_run = p.add_run()
                bullet_run.text = "• "
                bullet_run.font.size = Pt(14)
                bullet_run.font.name = "Calibri"
                bullet_run.font.color.rgb = RGBColor(0, 102, 204)

                # Add plain title with clickable "link" word
                if r_title and r_link:
                    # Add plain title
                    title_run = p.add_run()
                    title_run.text = r_title + " "
                    title_run.font.size = Pt(14)
                    title_run.font.name = "Calibri"
                    title_run.font.color.rgb = RGBColor(
                        0, 0, 0
                    )  # Black for better readability

                    # Add clickable "link" word
                    link_run = p.add_run()
                    link_run.text = "link"
                    link_run.font.size = Pt(14)
                    link_run.font.name = "Calibri"
                    link_run.font.color.rgb = RGBColor(0, 102, 204)
                    link_run.font.underline = True  # Make it look like a link

                    # Set hyperlink on the "link" word
                    try:
                        link_run.hyperlink.address = r_link
                    except Exception:
                        # If hyperlink setting fails, continue without it
                        pass

                elif r_title:
                    # Title without link
                    title_run = p.add_run()
                    title_run.text = r_title
                    title_run.font.size = Pt(14)
                    title_run.font.name = "Calibri"
                    title_run.font.color.rgb = RGBColor(
                        0, 0, 0
                    )  # Black for better readability

                elif r_link:
                    # Link without title - show clickable "link" word
                    link_run = p.add_run()
                    link_run.text = "link"
                    link_run.font.size = Pt(14)
                    link_run.font.name = "Calibri"
                    link_run.font.color.rgb = RGBColor(0, 102, 204)
                    link_run.font.underline = True

                    try:
                        link_run.hyperlink.address = r_link
                    except Exception:
                        pass

                else:
                    # Fallback: show string representation
                    fallback_run = p.add_run()
                    fallback_run.text = str(ref)
                    fallback_run.font.size = Pt(14)
                    fallback_run.font.name = "Calibri"
                    fallback_run.font.color.rgb = RGBColor(0, 102, 204)

            top += est_h + gap

        # Save PPT
        if in_memory:
            buffer = io.BytesIO()
            prs.save(buffer)
            data = buffer.getvalue()
            buffer.close()
            return data
        else:
            prs.save(output_filename)

    except Exception as e:
        print(f"Failed to generate PPT {output_filename}: {e}")


def estimate_height_wrapped_content(text, chars_per_line=100, line_height_pt=18):
    lines = 0
    for para in text.split("\n"):
        para = para.strip()
        if not para:
            continue
        lines += max(1, int(len(para) / chars_per_line) + 1)
    return Pt(lines * line_height_pt).inches


def estimate_height_wrapped_Title(text, chars_per_line=65, line_height_pt=20):
    lines = 0
    for para in text.split("\n"):
        para = para.strip()
        if not para:
            continue
        lines += max(1, int(len(para) / chars_per_line) + 1)
    return Pt(lines * line_height_pt).inches


def add_textbox(slide, title, content, top_inch):
    left = Inches(0.5)
    top = Inches(top_inch)
    width = Inches(9.5)
    height_inches = estimate_height_wrapped_content(content)
    height = Inches(height_inches)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    tf.clear()

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT

    run_title = p.add_run()
    run_title.text = f"{title}:\n"
    run_title.font.size = Pt(16)
    run_title.font.bold = True
    run_title.font.name = "Calibri"
    run_title.font.color.rgb = RGBColor(0x00, 0x66, 0xCC)

    run_content = p.add_run()
    run_content.text = content
    run_content.font.size = Pt(15)
    run_content.font.name = "Calibri"

    return top_inch + height_inches + gap


def add_textbox_Title(slide, title, content, top_inch):
    left = Inches(0.5)
    top = Inches(top_inch)
    width = Inches(9.5)
    height_inches = estimate_height_wrapped_Title(content)
    height = Inches(height_inches)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    tf.clear()

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT

    run_title = p.add_run()
    run_title.text = f"{title}: "
    run_title.font.size = Pt(20)
    run_title.font.name = "Calibri"
    run_title.font.bold = True
    run_title.font.color.rgb = RGBColor(0x00, 0x66, 0xCC)

    run_content = p.add_run()
    run_content.text = content
    run_content.font.size = Pt(20)
    run_content.font.name = "Calibri"
    run_content.font.bold = True
    run_content.font.color.rgb = RGBColor(0x00, 0x66, 0xCC)

    return top_inch + height_inches + gap
