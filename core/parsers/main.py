import pandas as pd
import uuid
import os
import shutil
from pathlib import Path
import asyncio
import fitz
import time
import markdown
from bs4 import BeautifulSoup
from PIL import Image
import io
import re
from app.socket_handler import sio
from core.parsers.image import image_parser
from core.models.document import Document, Page
from core.parsers.extensions import SUPPORTED_EXTENSIONS, IMAGE_EXTENSIONS
from pptx import Presentation
import traceback
import olefile

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_ROOT = os.path.dirname(BASE_DIR)


def extract_text_from_doc(path: str) -> str:
    """Extract readable text from a legacy .doc file (pure Python)."""
    if not olefile.isOleFile(path):
        raise ValueError(f"{path} is not a valid .doc file")

    with olefile.OleFileIO(path) as ole:
        if not ole.exists("WordDocument"):
            raise ValueError("No WordDocument stream found")
        stream = ole.openstream("WordDocument")
        data = stream.read()

    # Decode binary to text (best effort)
    text = data.decode("latin-1", errors="ignore")
    # Remove control characters
    text = re.sub(r"[\x00-\x08\x0B\x0C\x0E-\x1F]+", " ", text)
    # Collapse extra whitespace
    text = re.sub(r"\s{2,}", " ", text)
    # Keep only readable ASCII chunks
    text = "\n".join(re.findall(r"[ -~]{5,}", text))
    return text.strip()


async def extract_document(path, title="Untitled", file_name=None, thread_id=None):
    start_time = time.time()
    file_path = path
    ext = Path(path).suffix.lower()
    # Derive a safe base name even if file_name is None
    try:
        safe_file_name = file_name or os.path.basename(file_path)
        name, _ = os.path.splitext(safe_file_name)
    except Exception:
        traceback.print_exc()
        safe_file_name = os.path.basename(file_path)
        name, _ = os.path.splitext(safe_file_name)

    # Normalize thread to avoid crashing on None
    thread_id = thread_id or "unknown_thread"

    if ext not in SUPPORTED_EXTENSIONS:
        print(f"Unsupported file type: {ext} for {safe_file_name}. Skipping.")
        return None

    # --- Handle standalone images ---
    if ext in IMAGE_EXTENSIONS:
        try:
            text = await image_parser(file_path)
        except Exception as e:
            print(f"Error processing image {safe_file_name}: {str(e)}")
            traceback.print_exc()
            return None

        doc_id = str(uuid.uuid4())
        end_time = time.time()
        print(
            f"Time taken to process {safe_file_name} main image: {end_time - start_time} seconds"
        )
        return Document(
            id=doc_id,
            type=ext[1:],
            file_name=safe_file_name,
            content=[Page(number=1, text=text)],
            title=title,
            full_text=text,
        )

    # --- Handle Markdown files ---
    if ext == ".md":
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                md_text = f.read()

            # Convert markdown -> HTML -> plain text
            try:
                html = markdown.markdown(md_text)
            except Exception:
                traceback.print_exc()
                html = md_text  # fallback
            try:
                soup = BeautifulSoup(html, "html.parser")
                plain_text = soup.get_text(separator="\n")
            except Exception:
                traceback.print_exc()
                plain_text = md_text

            # Prepare image handling
            image_dir = f"data/threads/{thread_id}/images/{name}"
            try:
                os.makedirs(image_dir, exist_ok=True)
            except Exception:
                traceback.print_exc()

            ocr_tasks = {}
            image_names = []

            # Regex to find Markdown image syntax: ![alt](path)
            image_pattern = re.compile(r"!\[.*?\]\((.*?)\)")
            matches = image_pattern.findall(md_text)

            page_text = plain_text
            for idx, img_path in enumerate(matches, start=1):
                try:
                    resolved_path = img_path
                    if not os.path.isabs(resolved_path):
                        # make path relative to md file
                        resolved_path = os.path.join(
                            os.path.dirname(file_path), resolved_path
                        )

                    if not os.path.exists(resolved_path):
                        print(f"Markdown image not found: {resolved_path}")
                        continue

                    ext_img = Path(resolved_path).suffix.lstrip(".")
                    image_name = f"md_img{idx}.{ext_img}"
                    dest_path = os.path.join(image_dir, image_name)

                    # Copy image into project folder
                    try:
                        shutil.copy(resolved_path, dest_path)
                    except Exception:
                        traceback.print_exc()
                        continue
                    image_names.append(image_name)

                    placeholder = f"{{PENDING_{image_name}}}"
                    page_text += f"\n\n{placeholder}"

                    # Run OCR asynchronously
                    ocr_tasks[placeholder] = asyncio.create_task(
                        image_parser(dest_path)
                    )
                except Exception:
                    traceback.print_exc()

            # Wait for OCR tasks
            for placeholder, task in ocr_tasks.items():
                try:
                    image_text = await task
                except Exception as e:
                    print(f"Error parsing Markdown image: {e}")
                    traceback.print_exc()
                    image_text = "[Image OCR failed]"

                page_text = page_text.replace(placeholder, image_text, 1)

            doc_id = str(uuid.uuid4())

            return Document(
                id=doc_id,
                type="markdown",
                file_name=safe_file_name,
                content=[Page(number=1, text=page_text, images=image_names)],
                title=title,
                full_text=md_text,  # preserve original markdown
            )

        except Exception as e:
            print(f"Error processing Markdown file {safe_file_name}: {str(e)}")
            traceback.print_exc()
            return None

    if ext in {".xls", ".xlsx", ".csv"}:
        try:
            # Read Excel or CSV file
            if ext == ".xlsx":
                df = pd.read_excel(file_path, engine="openpyxl")
            elif ext == ".xls":
                df = pd.read_excel(file_path, engine="xlrd")
            else:
                df = pd.read_csv(file_path)

            # Clean merged cells & empty rows
            # df = df.ffill().dropna(how='all')

            # Normalize newlines inside cells
            df = df.applymap(
                lambda x: str(x).replace("\n", " ") if isinstance(x, str) else x
            )

            try:
                text = df.to_json(orient="records", lines=True)
                # text = df.to_markdown(index=False)

            except Exception:
                print("Error converting DataFrame to string")
                traceback.print_exc()
                text = str(df)

            # Optional: compact whitespace
            text = re.sub(r"\s{2,}", " ", text).strip()

            doc_id = str(uuid.uuid4())

            return Document(
                id=doc_id,
                type="spreadsheet",
                file_name=safe_file_name,
                content=[Page(number=1, text=text)],
                title=title,
                full_text=text,
            )

        except Exception as e:
            print(f"Error processing Excel/CSV file {safe_file_name}: {str(e)}")
            traceback.print_exc()
            return None

    # --- Handle legacy Word .doc files (single page, no image parsing) ---
    if ext == ".doc":
        try:
            text = extract_text_from_doc(file_path)
        except Exception as e:
            print(f"Error processing .doc file {safe_file_name}: {str(e)}")
            traceback.print_exc()
            return None

        doc_id = str(uuid.uuid4())
        end_time = time.time()
        print(
            f"Time taken to process {safe_file_name} (.doc): {end_time - start_time} seconds"
        )
        return Document(
            id=doc_id,
            type=ext[1:],
            file_name=safe_file_name,
            content=[Page(number=1, text=text)],
            title=title,
            full_text=text,
        )

    # --- Handle PowerPoint files ---
    if ext in {".ppt", ".pptx"}:
        try:
            prs = Presentation(file_path)
        except Exception as e:
            print(f"Error opening presentation {safe_file_name}: {e}")
            traceback.print_exc()
            return None
        pages = []
        combined_texts = []
        ocr_tasks = {}
        image_dir = f"data/threads/{thread_id}/images/{name}"
        try:
            os.makedirs(image_dir, exist_ok=True)
        except Exception:
            traceback.print_exc()

        try:
            for slide_number, slide in enumerate(prs.slides, start=1):
                # Extract text
                slide_text = []
                for shape in slide.shapes:
                    try:
                        if (
                            hasattr(shape, "text")
                            and getattr(shape, "text", "").strip()
                        ):
                            slide_text.append(shape.text.strip())
                    except Exception:
                        traceback.print_exc()
                page_text = "\n".join(slide_text)

                image_names = []

                # Extract images
                for shape_index, shape in enumerate(slide.shapes, start=1):
                    try:
                        if getattr(shape, "shape_type", None) == 13:  # PICTURE
                            image = shape.image
                            image_bytes = image.blob
                            image_ext = image.ext
                            image_name = (
                                f"slide{slide_number}_img{shape_index}.{image_ext}"
                            )
                            image_path = os.path.join(image_dir, image_name)

                            try:
                                with open(image_path, "wb") as f:
                                    f.write(image_bytes)
                            except Exception:
                                traceback.print_exc()
                                continue

                            placeholder = f"{{PENDING_{image_name}}}"
                            page_text += f"\n\n{placeholder}"
                            image_names.append(image_name)

                            ocr_tasks[placeholder] = asyncio.create_task(
                                image_parser(image_path)
                            )
                    except Exception:
                        traceback.print_exc()

                combined_texts.append(page_text)
                pages.append(
                    Page(number=slide_number, text=page_text, images=image_names)
                )

            # Wait for OCR tasks
            for placeholder, task in ocr_tasks.items():
                try:
                    image_text = await task
                except Exception as e:
                    print(f"Error parsing PPT image: {e}")
                    traceback.print_exc()
                    image_text = "[Image OCR failed]"

                for page in pages:
                    if placeholder in page.text:
                        page.text = page.text.replace(placeholder, image_text, 1)
                combined_texts = [
                    txt.replace(placeholder, image_text, 1) for txt in combined_texts
                ]
        except Exception:
            traceback.print_exc()

        doc_id = str(uuid.uuid4())
        end_time = time.time()
        print(
            f"Time taken to process {title} successfully: {end_time - start_time} seconds"
        )
        return Document(
            id=doc_id,
            type=ext[1:],
            file_name=safe_file_name,
            content=pages,
            title=title,
            full_text="\n".join(combined_texts),
        )

    # --- Handle PDFs ---
    if ext in [
        ".pdf",
        ".xlsx",
        ".epub",
        ".odt",
        ".txt",
        ".rtf",
        ".docx",
        ".html",
        ".xml",
    ]:
        try:
            doc = fitz.open(file_path)
        except Exception as e:
            print(f"Error opening PDF {safe_file_name}: {e}")
            traceback.print_exc()
            return None
        pages = []
        combined_texts = []
        ocr_tasks = {}

        image_dir_base = f"data/threads/{thread_id}/images/{name}"

        for page_number in range(len(doc)):
            try:
                page = doc.load_page(page_number)
                page_text = page.get_text("text")
            except Exception:
                traceback.print_exc()
                page_text = ""

            image_names = []
            image_dir = image_dir_base
            try:
                os.makedirs(image_dir, exist_ok=True)
            except Exception:
                traceback.print_exc()

            # Extract embedded raster images and schedule OCR only for these images
            try:
                image_list = page.get_images(full=True)
            except Exception:
                traceback.print_exc()
                image_list = []

            for img_index, img in enumerate(image_list):
                try:
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image.get("image")
                    image_ext = base_image.get("ext", "png")
                    if not image_bytes:
                        continue
                    image = Image.open(io.BytesIO(image_bytes))

                    image_name = f"page{page_number + 1}_img{img_index + 1}.{image_ext}"
                    image_path = os.path.join(image_dir, image_name)
                    try:
                        image.save(image_path)
                    except Exception:
                        traceback.print_exc()
                        continue

                    # Put placeholder where the image OCR result should go
                    placeholder = f"{{PENDING_{image_name}}}"
                    page_text += f"\n\n{placeholder}"
                    image_names.append(image_name)

                    # OCR only raster image files
                    ocr_tasks[placeholder] = asyncio.create_task(
                        image_parser(image_path)
                    )
                except Exception:
                    traceback.print_exc()

            combined_texts.append(page_text)
            pages.append(
                Page(number=page_number + 1, text=page_text, images=image_names)
            )

        # Wait for OCR tasks from the embedded raster images only
        for placeholder, task in ocr_tasks.items():
            try:
                image_text = await task
            except Exception as e:
                print(f"Error parsing image: {e}")
                traceback.print_exc()
                image_text = "[Image OCR failed]"

            # Replace placeholder once per occurrence (should be exactly 1)
            for page in pages:
                if placeholder in page.text:
                    page.text = page.text.replace(placeholder, image_text, 1)
            combined_texts = [
                txt.replace(placeholder, image_text, 1) for txt in combined_texts
            ]

        doc_id = str(uuid.uuid4())
        end_time = time.time()
        print(
            f"Time taken to process {title} successfully: {end_time - start_time} seconds"
        )
        return Document(
            id=doc_id,
            type=ext[1:],
            file_name=safe_file_name,
            content=pages,
            title=title,
            full_text="\n".join(combined_texts),
        )

    # If we reach here, the extension is supported but not yet implemented
    print(
        f"Parsing for file type {ext} not implemented for {safe_file_name}. Skipping."
    )
    return None
